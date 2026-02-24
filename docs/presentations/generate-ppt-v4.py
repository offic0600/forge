#!/usr/bin/env python3
"""
Forge Platform — Apple Keynote Style Executive Presentation v5
Design: Pure black, white typography, single accent, extreme restraint
Updated: 2026-02-24 — reflects feature-list v1.1 + planning baseline v2.2
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import lxml.etree as etree

# ── Apple-inspired palette (极度克制) ─────────────────────────
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TITLE = RGBColor(0xF5, 0xF5, 0xF7)    # Apple headline gray-white
GRAY_BODY = RGBColor(0x86, 0x86, 0x8B)      # Apple body gray
GRAY_DIM = RGBColor(0x48, 0x48, 0x4A)       # Dimmed text
GRAY_CARD = RGBColor(0x1D, 0x1D, 0x1F)      # Apple dark card
ACCENT = RGBColor(0x29, 0x97, 0xFF)         # Apple blue — 唯一强调色，少用

FONT = 'PingFang SC'
FONT_EN = 'SF Pro Display'
FONT_MONO = 'SF Mono'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def bg_black(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BLACK


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def rounded(slide, l, t, w, h, color, r=6000):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    sp = s._element
    prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is None:
            avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for child in list(avLst):
                avLst.remove(child)
        gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        gd.set('fmla', f'val {r}')
    return s


def txbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def txt(tf, s, sz=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font=FONT, lh=None):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if lh:
        p.line_spacing = Pt(lh)
    r = p.add_run()
    r.text = s
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def add(tf, s, sz=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font=FONT, before=6, after=2, lh=None):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if lh:
        p.line_spacing = Pt(lh)
    r = p.add_run()
    r.text = s
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def draw_logo(slide, cx, cy, scale=1.0):
    """
    Forge logo: abstract anvil + spark.
    """
    s = scale
    base = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(cx - 42*s), int(cy + 10*s),
        int(84*s), int(18*s)
    )
    base.fill.solid()
    base.fill.fore_color.rgb = GRAY_TITLE
    base.line.fill.background()

    neck = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(cx - 18*s), int(cy - 14*s),
        int(36*s), int(28*s)
    )
    neck.fill.solid()
    neck.fill.fore_color.rgb = GRAY_TITLE
    neck.line.fill.background()

    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(cx - 34*s), int(cy - 24*s),
        int(68*s), int(14*s)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = GRAY_TITLE
    top.line.fill.background()

    horn = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(cx + 34*s), int(cy - 20*s),
        int(22*s), int(8*s)
    )
    horn.fill.solid()
    horn.fill.fore_color.rgb = GRAY_TITLE
    horn.line.fill.background()

    spark = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        int(cx - 7*s), int(cy - 44*s),
        int(14*s), int(14*s)
    )
    spark.fill.solid()
    spark.fill.fore_color.rgb = ACCENT
    spark.line.fill.background()

    for dx, dy, sz_s in [(-18, -36, 5), (20, -38, 4), (-8, -52, 3), (14, -50, 3)]:
        p = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(cx + dx*s), int(cy + dy*s),
            int(sz_s*s), int(sz_s*s)
        )
        p.fill.solid()
        p.fill.fore_color.rgb = ACCENT
        p.line.fill.background()


def draw_logo_small(slide, cx, cy, scale=0.5):
    """Simplified small logo mark — just spark diamond."""
    spark = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        int(cx - 10*scale*2), int(cy - 10*scale*2),
        int(20*scale*2), int(20*scale*2)
    )
    spark.fill.solid()
    spark.fill.fore_color.rgb = ACCENT
    spark.line.fill.background()


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: Title — minimal, centered, logo
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s01_title():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    draw_logo(s, Inches(6.666), Inches(2.2), scale=Emu(Inches(0.012)))

    t = txbox(s, Inches(2), Inches(3.2), Inches(9.3), Inches(1.0))
    txt(t.text_frame, 'FORGE', sz=52, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER, font=FONT_EN)

    t2 = txbox(s, Inches(2), Inches(4.2), Inches(9.3), Inches(0.6))
    txt(t2.text_frame, '面向 AI 时代的智能交付平台', sz=22, color=GRAY_BODY, align=PP_ALIGN.CENTER)

    rect(s, Inches(5.8), Inches(5.1), Inches(1.7), Inches(0.02), GRAY_DIM)

    t3 = txbox(s, Inches(2), Inches(5.4), Inches(9.3), Inches(0.5))
    txt(t3.text_frame, '2026.02', sz=15, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: One line hook
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s02_hook():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5))
    txt(t.text_frame, '当 1 个人 + AI\n拥有 7 人团队的交付能力', sz=52, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER, lh=68)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3: The Problem — 4 pain points
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s03_problem():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '我们面对的现实', sz=40, bold=True, color=GRAY_TITLE)

    problems = [
        ('5-7 人的 Scrum 团队', '每个角色一个人，知识在人脑里，人走知识就没了'),
        ('AI 只是"更快的打字员"', '给每个人配一个 AI 助手，不等于交付效率的质变'),
        ('跨栈迁移是黑洞', '没人懂旧技术栈，人工迁移动辄数十人月'),
        ('经验不沉淀，错误在重复', '同样的部署失败、同样的安全漏洞，每个项目重来一遍'),
    ]

    for i, (title, desc) in enumerate(problems):
        y = Inches(2.0) + Inches(i * 1.25)

        rect(s, Inches(1.2), y + Inches(0.05), Inches(0.04), Inches(0.85), GRAY_DIM)

        t_title = txbox(s, Inches(1.6), y, Inches(10), Inches(0.45))
        txt(t_title.text_frame, title, sz=22, bold=True, color=GRAY_TITLE)

        t_desc = txbox(s, Inches(1.6), y + Inches(0.5), Inches(10), Inches(0.4))
        txt(t_desc.text_frame, desc, sz=16, color=GRAY_BODY)

    t5 = txbox(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.5))
    txt(t5.text_frame, '问题的本质不是工具不够，是交付流程本身没有被重构。', sz=18, color=ACCENT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4: Three Insights — updated numbers + learning loop
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s04_insight():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '三个关键洞察', sz=40, bold=True, color=GRAY_TITLE)

    insights = [
        (
            'SuperAgent\n不是 Multi-Agent',
            '不是 6 个角色各配一个 AI。\n而是 1 个超级智能体，\n通过 Skill Profile 动态切换。',
            '1 个智能体 > 6 个独立 Agent',
        ),
        (
            'Skill + 私域知识\n双护城河',
            '模型会被追平，Prompt 会被抄走。\n但 32 个 Skill 编码的工程经验\n+ 持续沉淀的私域知识，无法复制。',
            'Skill × Knowledge = 不可复制的数字资产',
        ),
        (
            '双环驱动\n越用越好',
            '交付环做事，进化环沉淀知识。\n每次交付产生私域数据，\n反哺下次交付——飞轮自动转。',
            '交付 × 知识沉淀 = 飞轮效应',
        ),
    ]

    for i, (title, desc, tagline) in enumerate(insights):
        x = Inches(0.8) + Inches(i * 4.15)
        y = Inches(1.9)

        c = rounded(s, x, y, Inches(3.85), Inches(4.6), GRAY_CARD)

        rect(s, x, y, Inches(3.85), Inches(0.03), ACCENT if i == 0 else GRAY_DIM)

        t_title = txbox(s, x + Inches(0.4), y + Inches(0.5), Inches(3.05), Inches(1.2))
        txt(t_title.text_frame, title, sz=24, bold=True, color=GRAY_TITLE, lh=32)

        t_desc = txbox(s, x + Inches(0.4), y + Inches(1.8), Inches(3.05), Inches(1.8))
        txt(t_desc.text_frame, desc, sz=15, color=GRAY_BODY, lh=24)

        t_tag = txbox(s, x + Inches(0.4), y + Inches(3.9), Inches(3.05), Inches(0.4))
        txt(t_tag.text_frame, tagline, sz=13, color=GRAY_DIM)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5: SuperAgent — 6 profiles + OODA
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s05_superagent():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '一个智能体，六种角色', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.3), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '根据任务自动切换 Skill Profile，覆盖从需求到运维再到评估的完整交付链', sz=16, color=GRAY_BODY)

    profiles = [
        ('规划', 'Planning', '需求分析\nPRD 编写'),
        ('设计', 'Design', '架构设计\nAPI 设计'),
        ('开发', 'Development', '代码生成\n32 Skills'),
        ('测试', 'Testing', '用例设计\n自动化执行'),
        ('运维', 'Ops', '部署策略\nK8s 编排'),
        ('评估', 'Evaluation', '进度评估\n知识蒸馏'),
    ]

    for i, (cn, en, desc) in enumerate(profiles):
        x = Inches(0.5) + Inches(i * 2.1)
        y = Inches(2.2)

        c = rounded(s, x, y, Inches(1.9), Inches(3.0), GRAY_CARD)
        # Development gets accent, evaluation gets accent too
        rect(s, x, y, Inches(1.9), Inches(0.03), ACCENT if i in (2, 5) else GRAY_DIM)

        t_cn = txbox(s, x + Inches(0.15), y + Inches(0.35), Inches(1.6), Inches(0.5))
        txt(t_cn.text_frame, cn, sz=24, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER)

        t_en = txbox(s, x + Inches(0.15), y + Inches(0.85), Inches(1.6), Inches(0.35))
        txt(t_en.text_frame, en, sz=11, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_MONO)

        t_desc = txbox(s, x + Inches(0.15), y + Inches(1.5), Inches(1.6), Inches(1.2))
        txt(t_desc.text_frame, desc, sz=13, color=GRAY_BODY, align=PP_ALIGN.CENTER, lh=20)

    # OODA bar
    ooda_y = Inches(5.7)
    rounded(s, Inches(0.8), ooda_y, Inches(11.7), Inches(1.2), GRAY_CARD)

    t_ooda_label = txbox(s, Inches(1.2), ooda_y + Inches(0.1), Inches(3), Inches(0.35))
    txt(t_ooda_label.text_frame, 'OODA 循环', sz=12, color=GRAY_DIM, font=FONT_EN)

    steps = ['Observe  观察', 'Orient  分析', 'Decide  决策', 'Act  执行', 'HITL  人在回路']
    for i, step in enumerate(steps):
        sx = Inches(1.0) + Inches(i * 2.35)
        st = txbox(s, sx, ooda_y + Inches(0.45), Inches(2.1), Inches(0.5))
        txt(st.text_frame, step, sz=15, color=GRAY_TITLE if i < 4 else ACCENT, align=PP_ALIGN.CENTER)

        if i < 4:
            arrow = txbox(s, sx + Inches(1.9), ooda_y + Inches(0.45), Inches(0.5), Inches(0.5))
            txt(arrow.text_frame, '→', sz=16, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6: Agentic Loop — 自主执行管线 (NEW)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s06_agentic_loop():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    txt(t.text_frame, '不是聊天，是自主执行', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.1), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '用户发一条消息，SuperAgent 自动启动 50 轮执行管线——无需逐步指令', sz=16, color=GRAY_BODY)

    # Pipeline steps — 7 steps in a vertical flow, left side
    pipeline = [
        ('1', 'Profile 路由', '分析意图，自动选择专业角色'),
        ('2', 'Skill 加载', '按 Profile + 关键词动态加载 3-20 个 Skill'),
        ('3', '记忆注入', 'Workspace Memory + Stage Memory + Session History'),
        ('4', 'OODA 循环', '观察→分析→决策→执行，调用 18 种 MCP 工具'),
        ('5', 'Context 管理', '每轮估算 Token 使用率，超限自动压缩'),
        ('6', '底线检查', '代码生成后自动运行质量底线，失败自修复'),
        ('7', '交付 + 知识沉淀', '写入文件 + 更新记忆 + 生成 Summary + 私域知识自动积累'),
    ]

    for i, (num, title, desc) in enumerate(pipeline):
        y = Inches(1.75) + Inches(i * 0.74)
        x = Inches(0.8)

        # Number circle
        t_num = txbox(s, x, y + Inches(0.05), Inches(0.45), Inches(0.45))
        txt(t_num.text_frame, num, sz=16, bold=True,
            color=ACCENT if i == 3 else GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

        # Title
        t_title = txbox(s, x + Inches(0.55), y, Inches(2.2), Inches(0.35))
        txt(t_title.text_frame, title, sz=17, bold=True,
            color=ACCENT if i == 3 else GRAY_TITLE)

        # Description
        t_desc = txbox(s, x + Inches(0.55), y + Inches(0.32), Inches(4.5), Inches(0.3))
        txt(t_desc.text_frame, desc, sz=12, color=GRAY_BODY)

        # Connecting line
        if i < 6:
            rect(s, x + Inches(0.2), y + Inches(0.52), Inches(0.02), Inches(0.22), GRAY_DIM)

    # Right side: 5 key capabilities
    cap_x = Inches(6.5)
    rounded(s, cap_x, Inches(1.75), Inches(6.0), Inches(5.4), GRAY_CARD)
    rect(s, cap_x, Inches(1.75), Inches(6.0), Inches(0.03), ACCENT)

    t_cap_title = txbox(s, cap_x + Inches(0.4), Inches(2.0), Inches(5.2), Inches(0.5))
    txt(t_cap_title.text_frame, '关键能力', sz=20, bold=True, color=ACCENT)

    capabilities = [
        ('自主多轮执行', '自行决定调用什么工具、读什么文件、写什么代码\n单次对话最多 50 轮'),
        ('工具链编排', '搜索知识→读取文件→生成代码→写入 Workspace\n→编译→运行测试，全自动串联'),
        ('失败自修复', '底线检查失败→分析原因→修改代码→重新检查\n最多 2 轮自动修复'),
        ('上下文自管理', 'Token 接近上限时自动压缩历史\n不中断执行（3 阶段渐进压缩）'),
        ('记忆增强', '每次执行携带项目记忆\n技术栈、关键决策、进行中的工作——无需重复说明'),
    ]

    for i, (cap_title, cap_desc) in enumerate(capabilities):
        cy = Inches(2.65) + Inches(i * 0.88)

        t_ct = txbox(s, cap_x + Inches(0.4), cy, Inches(5.2), Inches(0.3))
        txt(t_ct.text_frame, cap_title, sz=16, bold=True, color=GRAY_TITLE)

        t_cd = txbox(s, cap_x + Inches(0.4), cy + Inches(0.3), Inches(5.2), Inches(0.5))
        txt(t_cd.text_frame, cap_desc, sz=12, color=GRAY_BODY, lh=18)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7: Four-Layer Architecture — updated numbers
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s07_architecture():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    txt(t.text_frame, '四层解耦，独立演进', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.1), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '每一层可以独立升级，不影响其他层 — 这是平台能活过模型迭代的根本原因', sz=16, color=GRAY_BODY)

    layers = [
        (
            '用户交互层',
            'User Interaction',
            'CLI / VS Code (开发者) + Web IDE (全角色)',
            'Monaco 编辑器  ·  AI Chat  ·  知识浏览  ·  可视化工作流',
            False,
        ),
        (
            'SuperAgent 层',
            'Intelligence',
            '6 Profiles × 32 Skills × 8 Baselines',
            'Skill Profile Router  ·  Baseline Runner  ·  Agentic Loop 50 轮  ·  3 层记忆',
            True,
        ),
        (
            'Agent Runtime 层',
            'Runtime',
            'MCP 协议  ·  Tool Calling  ·  沙箱  ·  可观测',
            '18 聚合工具  ·  安全围栏  ·  SuperAgent 评估框架  ·  四维评分',
            False,
        ),
        (
            '数据与模型层',
            'Data & Model',
            '6 Provider + 私域知识库（三层 Scope：Global / Workspace / Personal）',
            'MCP Gateway → Knowledge / Database / ServiceGraph / 执行日志 → 知识沉淀闭环',
            False,
        ),
    ]

    for i, (cn, en, line1, line2, is_accent) in enumerate(layers):
        y = Inches(1.75) + Inches(i * 1.35)

        c = rounded(s, Inches(0.8), y, Inches(11.7), Inches(1.2), GRAY_CARD)

        rect(s, Inches(0.8), y, Inches(0.05), Inches(1.2), ACCENT if is_accent else GRAY_DIM)

        t_num = txbox(s, Inches(1.15), y + Inches(0.15), Inches(0.65), Inches(0.55))
        txt(t_num.text_frame, f'L{i+1}', sz=14, bold=True,
            color=ACCENT if is_accent else GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

        t_cn = txbox(s, Inches(1.9), y + Inches(0.1), Inches(2.5), Inches(0.4))
        txt(t_cn.text_frame, cn, sz=20, bold=True, color=GRAY_TITLE)

        t_en = txbox(s, Inches(4.3), y + Inches(0.15), Inches(2), Inches(0.35))
        txt(t_en.text_frame, en, sz=11, color=GRAY_DIM, font=FONT_MONO)

        t_l1 = txbox(s, Inches(1.9), y + Inches(0.55), Inches(10), Inches(0.3))
        txt(t_l1.text_frame, line1, sz=14, color=ACCENT if is_accent else GRAY_BODY)

        t_l2 = txbox(s, Inches(1.9), y + Inches(0.85), Inches(10), Inches(0.3))
        txt(t_l2.text_frame, line2, sz=12, color=GRAY_DIM)

        if i < 3:
            arr_y = y + Inches(1.2)
            t_arr = txbox(s, Inches(6.2), arr_y, Inches(0.9), Inches(0.2))
            txt(t_arr.text_frame, '↓', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

    t_bot = txbox(s, Inches(1.2), Inches(7.0), Inches(11), Inches(0.4))
    txt(t_bot.text_frame, '模型从 Claude 换成 GPT?  换 L4 的 Adapter。32 个 Skills、私域知识库和 8 个 Baselines 一行不改。', sz=14, color=ACCENT, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8: Stable vs Volatile — updated numbers
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s08_stable_volatile():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    txt(t.text_frame, '稳态固化 · 敏态抽象', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.1), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '模型一定会被追平，Skill + 私域知识才是真正的双护城河', sz=16, color=GRAY_BODY)

    # ── Left: 稳态 ──
    rounded(s, Inches(0.8), Inches(1.8), Inches(4.8), Inches(4.8), GRAY_CARD)
    rect(s, Inches(0.8), Inches(1.8), Inches(4.8), Inches(0.03), ACCENT)

    t_sl = txbox(s, Inches(1.2), Inches(2.1), Inches(4), Inches(0.5))
    txt(t_sl.text_frame, '稳态 — 核心资产', sz=20, bold=True, color=ACCENT)

    t_sld = txbox(s, Inches(1.2), Inches(2.6), Inches(4), Inches(0.3))
    txt(t_sld.text_frame, '越积越厚，无法复制', sz=13, color=GRAY_DIM)

    stable_items = [
        ('32 Skills', '十年工程经验的数字化编码'),
        ('私域知识库', '三层 Scope · 13+ 文档 · 持续沉淀'),
        ('8 Baselines', '质量底线自动守护'),
        ('6 Profiles', '全交付角色覆盖 + 评估'),
        ('OODA 流程', '标准化决策环路'),
    ]

    for i, (item, desc) in enumerate(stable_items):
        iy = Inches(3.1) + Inches(i * 0.64)
        t_item = txbox(s, Inches(1.5), iy, Inches(2.2), Inches(0.35))
        txt(t_item.text_frame, item, sz=17, bold=True, color=GRAY_TITLE, font=FONT_EN)
        t_desc = txbox(s, Inches(3.6), iy + Inches(0.02), Inches(2), Inches(0.35))
        txt(t_desc.text_frame, desc, sz=13, color=GRAY_BODY)

    # ── Center: Adapter ──
    mid_x = Inches(5.95)
    mid_y = Inches(3.2)

    adapter = rounded(s, mid_x - Inches(0.55), mid_y, Inches(1.1), Inches(2.0), GRAY_DIM, r=4000)

    t_ad = txbox(s, mid_x - Inches(0.55), mid_y + Inches(0.2), Inches(1.1), Inches(1.6))
    tf = t_ad.text_frame
    tf.word_wrap = True
    txt(tf, 'A', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN)
    add(tf, 'd', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 'a', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 'p', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 't', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 'e', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 'r', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)

    t_al = txbox(s, Inches(5.6), Inches(3.9), Inches(0.5), Inches(0.4))
    txt(t_al.text_frame, '→', sz=18, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)

    t_ar = txbox(s, Inches(6.9), Inches(3.9), Inches(0.5), Inches(0.4))
    txt(t_ar.text_frame, '→', sz=18, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

    # ── Right: 敏态 ──
    rounded(s, Inches(7.7), Inches(1.8), Inches(4.8), Inches(4.8), GRAY_CARD)
    rect(s, Inches(7.7), Inches(1.8), Inches(4.8), Inches(0.03), GRAY_DIM)

    t_vl = txbox(s, Inches(8.1), Inches(2.1), Inches(4), Inches(0.5))
    txt(t_vl.text_frame, '敏态 — 可替换组件', sz=20, bold=True, color=GRAY_BODY)

    t_vld = txbox(s, Inches(8.1), Inches(2.6), Inches(4), Inches(0.3))
    txt(t_vld.text_frame, '技术迭代时只换这里', sz=13, color=GRAY_DIM)

    volatile_items = [
        ('ModelAdapter', '6 Provider 已就绪\nClaude/Gemini/Qwen/Bedrock/OpenAI/MiniMax'),
        ('RuntimeAdapter', 'Claude Code → ForgeNative'),
        ('AssetFormatAdapter', 'SKILL.md v1 → v2 → ...'),
    ]

    for i, (item, desc) in enumerate(volatile_items):
        iy = Inches(3.1) + Inches(i * 0.8)
        t_item = txbox(s, Inches(8.3), iy, Inches(4), Inches(0.35))
        txt(t_item.text_frame, item, sz=16, bold=True, color=GRAY_TITLE, font=FONT_MONO)
        t_desc = txbox(s, Inches(8.3), iy + Inches(0.35), Inches(4), Inches(0.3))
        txt(t_desc.text_frame, desc, sz=13, color=GRAY_DIM, lh=18)

    # Adapter labels
    adapters_label = [
        'ModelAdapter.kt  ·  RuntimeAdapter.kt  ·  AssetFormatAdapter.kt'
    ]
    t_alab = txbox(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.3))
    txt(t_alab.text_frame, adapters_label[0], sz=12, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_MONO)

    # Bottom insight
    rounded(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.85), GRAY_CARD)
    rect(s, Inches(0.8), Inches(6.2), Inches(0.05), Inches(0.85), ACCENT)

    t_bot = txbox(s, Inches(1.3), Inches(6.35), Inches(10.8), Inches(0.6))
    txt(t_bot.text_frame, '战略意义', sz=13, bold=True, color=ACCENT)
    add(t_bot.text_frame,
        'OpenAI 发布更好的模型？换一个 Adapter。6 家模型已热插拔就绪。\n'
        '但 32 个 Skill + 持续沉淀的私域知识——这套数字资产，竞争对手无法复制。',
        sz=14, color=GRAY_BODY, before=4, lh=22)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9: The Breakthrough — Before / After (updated)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s09_breakthrough():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, 'AI 不再只是聊天，而是交付', sz=40, bold=True, color=GRAY_TITLE)

    # BEFORE
    rounded(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), GRAY_CARD)
    rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.03), GRAY_DIM)

    t_b = txbox(s, Inches(1.2), Inches(2.1), Inches(4.8), Inches(0.5))
    txt(t_b.text_frame, 'Before', sz=18, bold=True, color=GRAY_DIM, font=FONT_EN)

    before = [
        'AI 在聊天窗口展示代码',
        '用户手动复制粘贴到文件',
        '每次对话从零开始，知识不沉淀',
        '团队经验在人脑里，人走就没了',
        '知识搜索返回 mock 数据',
        'AI 不了解项目结构',
    ]
    t_items = txbox(s, Inches(1.2), Inches(2.7), Inches(4.8), Inches(3.5))
    txt(t_items.text_frame, '', sz=8)
    for item in before:
        add(t_items.text_frame, f'  {item}', sz=16, color=GRAY_BODY, before=10, after=2)

    # Arrow
    a = txbox(s, Inches(6.2), Inches(3.6), Inches(0.9), Inches(0.8))
    txt(a.text_frame, '→', sz=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)

    # AFTER
    rounded(s, Inches(6.9), Inches(1.8), Inches(5.9), Inches(4.8), GRAY_CARD)
    rect(s, Inches(6.9), Inches(1.8), Inches(5.9), Inches(0.03), ACCENT)

    t_a = txbox(s, Inches(7.3), Inches(2.1), Inches(5.1), Inches(0.5))
    txt(t_a.text_frame, 'After — Forge', sz=18, bold=True, color=ACCENT, font=FONT_EN)

    after = [
        ('Agent 自主执行 50 轮', 'Agentic Loop 自动规划→编码→验证'),
        ('AI 直接写文件到 workspace', 'workspace_write + SSE file_changed'),
        ('3 层记忆跨 Session 携带', 'Workspace + Stage + Session Memory'),
        ('交付即沉淀——私域知识自动积累', 'search + page_create + 三层 Scope'),
        ('四维评估 + 学习闭环', '每次交付自动评分 → 反哺下次交付'),
        ('6 家模型自由切换 + SSO', 'Claude/Gemini/Qwen/Bedrock/OpenAI/MiniMax'),
    ]
    t_items2 = txbox(s, Inches(7.3), Inches(2.7), Inches(5.1), Inches(3.5))
    txt(t_items2.text_frame, '', sz=8)
    for item, note in after:
        add(t_items2.text_frame, f'  {item}', sz=16, color=GRAY_TITLE, before=6, after=0)
        add(t_items2.text_frame, f'     {note}', sz=11, color=GRAY_DIM, before=0, after=4, font=FONT_MONO)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10: Numbers — updated metrics
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s10_numbers():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '我们走到了哪里', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.25), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, 'Phase 0 → 7  全部完成（8 个里程碑，32 个 Session）', sz=16, color=GRAY_BODY)

    metrics = [
        ('50,000+', '行代码'),
        ('68', 'REST API'),
        ('156', '单元测试'),
        ('38', '验收用例'),
    ]

    for i, (num, label) in enumerate(metrics):
        x = Inches(0.8) + Inches(i * 3.15)
        y = Inches(2.4)

        t_num = txbox(s, x, y, Inches(2.9), Inches(1.5))
        txt(t_num.text_frame, num, sz=64, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER, font=FONT_EN)

        t_label = txbox(s, x, y + Inches(1.3), Inches(2.9), Inches(0.4))
        txt(t_label.text_frame, label, sz=16, color=GRAY_BODY, align=PP_ALIGN.CENTER)

    metrics2 = [
        ('18', 'MCP 工具'),
        ('6', 'Skill Profile'),
        ('32', 'Skills'),
        ('6', 'Docker 容器'),
    ]

    for i, (num, label) in enumerate(metrics2):
        x = Inches(0.8) + Inches(i * 3.15)
        y = Inches(4.5)

        t_num = txbox(s, x, y, Inches(2.9), Inches(1.3))
        txt(t_num.text_frame, num, sz=56, bold=True, color=ACCENT if i == 0 else GRAY_TITLE, align=PP_ALIGN.CENTER, font=FONT_EN)

        t_label = txbox(s, x, y + Inches(1.1), Inches(2.9), Inches(0.4))
        txt(t_label.text_frame, label, sz=16, color=GRAY_BODY, align=PP_ALIGN.CENTER)

    rect(s, Inches(1.2), Inches(4.3), Inches(10.9), Inches(0.01), GRAY_DIM)

    t_bot = txbox(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4))
    txt(t_bot.text_frame, 'docker compose up → SSO → Agent 50 轮自主执行 → 代码写入 workspace → 私域知识沉淀 → 学习闭环', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11: Cross-stack Migration — kept
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s11_migration():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '跨栈迁移：10 倍提效', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.3), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '"没人懂旧技术栈"的困局，AI 可以系统性地解决', sz=16, color=GRAY_BODY)

    steps = [
        ('01', 'AI 代码考古', '解析项目结构\n生成项目说明书'),
        ('02', '业务规则提取', '逐模块分析\n规则 + 边界条件'),
        ('03', '架构设计', '.NET → Java 映射\n依赖拓扑排列'),
        ('04', '逐模块迁移', 'SuperAgent 编写\n规范自动遵循'),
        ('05', '行为等价验证', 'API 契约对比\n规则覆盖度检查'),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.5) + Inches(i * 2.52)
        y = Inches(2.2)

        c = rounded(s, x, y, Inches(2.3), Inches(2.6), GRAY_CARD)
        rect(s, x, y, Inches(2.3), Inches(0.03), ACCENT if i == 0 else GRAY_DIM)

        t_num = txbox(s, x + Inches(0.2), y + Inches(0.25), Inches(0.5), Inches(0.35))
        txt(t_num.text_frame, num, sz=13, color=GRAY_DIM, font=FONT_EN)

        t_title = txbox(s, x + Inches(0.2), y + Inches(0.65), Inches(1.9), Inches(0.5))
        txt(t_title.text_frame, title, sz=17, bold=True, color=GRAY_TITLE)

        t_desc = txbox(s, x + Inches(0.2), y + Inches(1.2), Inches(1.9), Inches(1.2))
        txt(t_desc.text_frame, desc, sz=13, color=GRAY_BODY, lh=20)

        if i < 4:
            arrow = txbox(s, x + Inches(2.15), y + Inches(0.8), Inches(0.55), Inches(0.5))
            txt(arrow.text_frame, '→', sz=18, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Impact comparison
    y2 = Inches(5.2)
    rect(s, Inches(1.2), y2, Inches(10.9), Inches(0.01), GRAY_DIM)

    t_trad_l = txbox(s, Inches(1.5), y2 + Inches(0.3), Inches(3), Inches(0.3))
    txt(t_trad_l.text_frame, '传统方式', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER)
    t_trad_n = txbox(s, Inches(1.5), y2 + Inches(0.6), Inches(3), Inches(0.8))
    txt(t_trad_n.text_frame, '50-65 人天', sz=32, bold=True, color=GRAY_BODY, align=PP_ALIGN.CENTER, font=FONT_EN)

    t_arr = txbox(s, Inches(5), y2 + Inches(0.6), Inches(1), Inches(0.6))
    txt(t_arr.text_frame, '→', sz=28, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

    t_forge_l = txbox(s, Inches(6.2), y2 + Inches(0.3), Inches(3), Inches(0.3))
    txt(t_forge_l.text_frame, 'Forge 辅助', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER)
    t_forge_n = txbox(s, Inches(6.2), y2 + Inches(0.6), Inches(3), Inches(0.8))
    txt(t_forge_n.text_frame, '5-7 人天', sz=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)

    t_mult = txbox(s, Inches(9.8), y2 + Inches(0.4), Inches(2.5), Inches(1.0))
    txt(t_mult.text_frame, '10x', sz=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12: Quality — the safety net story
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s12_quality():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '速度不缺，缺的是安全网', sz=40, bold=True, color=GRAY_TITLE)

    t_story = txbox(s, Inches(1.2), Inches(1.8), Inches(10), Inches(1.8))
    txt(t_story.text_frame, '首次部署成功率', sz=14, color=GRAY_DIM)
    add(t_story.text_frame, '7.7%', sz=72, bold=True, color=GRAY_TITLE, font=FONT_EN, before=0, after=0)
    add(t_story.text_frame, '13 次尝试，90 分钟，10 个计划外问题。CORS 通配符直接进入 commit，无人察觉。', sz=16, color=GRAY_BODY, before=8)

    rect(s, Inches(1.2), Inches(4.0), Inches(10.9), Inches(0.01), GRAY_DIM)

    t_sol = txbox(s, Inches(1.2), Inches(4.3), Inches(5.5), Inches(1.8))
    txt(t_sol.text_frame, '系统性解决', sz=20, bold=True, color=GRAY_TITLE)
    add(t_sol.text_frame, '痛点不修补，而是编码为平台能力。', sz=16, color=GRAY_BODY, before=10)
    add(t_sol.text_frame, '8 道质量底线自动拦截 + 156 个单元测试，', sz=16, color=GRAY_BODY, before=4)
    add(t_sol.text_frame, '让未来的开发者无感地获得保护。', sz=16, color=GRAY_BODY, before=4)

    t_target = txbox(s, Inches(7.5), Inches(4.3), Inches(5), Inches(1.8))
    txt(t_target.text_frame, '验收测试通过率', sz=14, color=GRAY_DIM)
    add(t_target.text_frame, '94.7%', sz=56, bold=True, color=ACCENT, font=FONT_EN, before=0)
    add(t_target.text_frame, '38 TC / 11 场景', sz=14, color=GRAY_BODY, before=4)

    t_p = txbox(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4))
    txt(t_p.text_frame, '已验证的设计不可隐式退化  ·  开发者痛点即平台能力  ·  四维评估持续监控', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13: Roadmap — completely updated
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s13_roadmap():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '从工具到平台到生态', sz=40, bold=True, color=GRAY_TITLE)

    cols = [
        ('已完成 ✅', ACCENT, [
            ('Phase 0-1', '骨架 + Web IDE 实连'),
            ('Phase 1.5-1.6', '设计守护 + Docker + AI 交付闭环'),
            ('Phase 2-3', '多模型 + HITL + 质量面板'),
            ('Phase 4-5', 'Skill 架构 + 3 层记忆'),
            ('Phase 6-7', '产品加固 + 评估 + 学习闭环框架'),
        ]),
        ('下一步 — 知识链路闭合', GRAY_TITLE, [
            ('Phase 8', '进化环全面打通'),
            ('', '知识自动提取 → 私域知识库沉淀'),
            ('', 'Convention 挖掘 → Skill 自动更新'),
            ('', '知识空白检测 → 自动生成文档'),
            ('', '执行日志 → 四维评估 → 反哺交付'),
        ]),
        ('远景', GRAY_DIM, [
            ('Phase 9', '生产就绪 + 团队协作'),
            ('', 'PostgreSQL + 多租户 + CI/CD'),
            ('Phase 10', '生态扩展'),
            ('', 'Skill Marketplace + CLI'),
        ]),
    ]

    for ci, (section, color, items) in enumerate(cols):
        x = Inches(0.8) + Inches(ci * 4.15)
        y = Inches(1.8)

        t_sec = txbox(s, x, y, Inches(3.85), Inches(0.5))
        txt(t_sec.text_frame, section, sz=14, bold=True, color=color, font=FONT_EN)

        rect(s, x, y + Inches(0.5), Inches(3.85), Inches(0.02), color)

        c = rounded(s, x, y + Inches(0.7), Inches(3.85), Inches(4.3), GRAY_CARD)

        t_items = txbox(s, x + Inches(0.35), y + Inches(1.0), Inches(3.15), Inches(3.8))
        txt(t_items.text_frame, '', sz=8)

        for phase, desc in items:
            if phase:
                add(t_items.text_frame, phase, sz=15, bold=True, color=GRAY_TITLE, before=14, after=2)
                add(t_items.text_frame, desc, sz=14, color=GRAY_BODY, before=2, after=2)
            else:
                add(t_items.text_frame, desc, sz=14, color=GRAY_BODY, before=6, after=2)

    t_v = txbox(s, Inches(1.2), Inches(6.6), Inches(11), Inches(0.4))
    txt(t_v.text_frame, '终局：Skill + 私域知识双飞轮  ·  交付即沉淀  ·  越用越强', sz=15, color=GRAY_DIM, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14: The Ask — updated
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s14_ask():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '下一步需要什么', sz=40, bold=True, color=GRAY_TITLE)

    asks = [
        ('01', '1 个试点项目', '选一个真实业务项目（迁移或新建均可），验证 Skill + 知识 双飞轮。\n全流程跑通：交付 → 知识沉淀 → 下次交付更快——验证进化环价值。'),
        ('02', '5-10 人内部试用', 'Web IDE 已支持 SSO + 多 Workspace + 6 家模型。\n真实用户的交付数据是知识链路闭合最关键的燃料。'),
        ('03', '知识运营支持', '知识管理员 + 领域专家协助梳理私域知识初始集。\n框架已就绪，需要真实的业务知识灌入，让飞轮转起来。'),
    ]

    for i, (num, title, desc) in enumerate(asks):
        y = Inches(1.9) + Inches(i * 1.7)

        rounded(s, Inches(0.8), y, Inches(11.7), Inches(1.4), GRAY_CARD)
        rect(s, Inches(0.8), y, Inches(0.05), Inches(1.4), ACCENT if i == 0 else GRAY_DIM)

        t_num = txbox(s, Inches(1.2), y + Inches(0.25), Inches(0.8), Inches(0.5))
        txt(t_num.text_frame, num, sz=28, bold=True, color=ACCENT if i == 0 else GRAY_DIM, font=FONT_EN)

        t_title = txbox(s, Inches(2.1), y + Inches(0.15), Inches(5), Inches(0.5))
        txt(t_title.text_frame, title, sz=24, bold=True, color=GRAY_TITLE)

        t_desc = txbox(s, Inches(2.1), y + Inches(0.65), Inches(9), Inches(0.6))
        txt(t_desc.text_frame, desc, sz=14, color=GRAY_BODY, lh=22)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 15: Closing — updated tagline
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s15_closing():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    draw_logo(s, Inches(6.666), Inches(1.8), scale=Emu(Inches(0.008)))

    t1 = txbox(s, Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.8))
    txt(t1.text_frame, '每一次交付，都在沉淀知识', sz=28, color=GRAY_DIM, align=PP_ALIGN.CENTER)

    t2 = txbox(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.0))
    txt(t2.text_frame, '每一次沉淀，都让下次交付更强', sz=42, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER)

    t_sub2 = txbox(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.5))
    txt(t_sub2.text_frame, 'Skill + 私域知识 — 不可复制的双护城河', sz=18, color=ACCENT, align=PP_ALIGN.CENTER)

    rect(s, Inches(5.8), Inches(5.1), Inches(1.7), Inches(0.02), GRAY_DIM)

    t3 = txbox(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5))
    txt(t3.text_frame, 'Thank You', sz=20, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)


# ── Generate ──────────────────────────────────────────────────
s01_title()
s02_hook()
s03_problem()
s04_insight()
s05_superagent()
s06_agentic_loop()
s07_architecture()
s08_stable_volatile()
s09_breakthrough()
s10_numbers()
s11_migration()
s12_quality()
s13_roadmap()
s14_ask()
s15_closing()

out = 'docs/presentations/forge-platform-executive-v5.pptx'
prs.save(out)
print(f'✅ {out} ({len(prs.slides)} slides)')
